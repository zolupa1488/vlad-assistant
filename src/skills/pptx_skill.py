"""Generate a clean .pptx presentation from a structured outline.

Produces a deck that doesn't look like a 2003 PowerPoint default — sane fonts,
proper margins, accent stripe on title slide, simple bullets on content slides.
Designed to be called from Claude as a tool, with the outline produced by the
LLM itself.
"""

from __future__ import annotations

import os
import uuid
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

OUTPUT_DIR = os.environ.get("FILES_OUTPUT_DIR", "/app/data/files")

# Subtle palette — close to claude.ai
_BG = RGBColor(0xF8, 0xF8, 0xF6)
_INK = RGBColor(0x1F, 0x1F, 0x1F)
_DIM = RGBColor(0x55, 0x55, 0x55)
_ACCENT = RGBColor(0xCC, 0x78, 0x5C)


def _add_textbox(slide, *, x: float, y: float, w: float, h: float, text: str,
                 size_pt: int, bold: bool = False, color: RGBColor | None = None) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.name = "Helvetica"
    run.font.color.rgb = color or _INK


def _add_bullets(slide, *, x: float, y: float, w: float, h: float, items: list[str],
                 size_pt: int = 18) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f"•  {item}"
        run.font.size = Pt(size_pt)
        run.font.name = "Helvetica"
        run.font.color.rgb = _INK


def _paint_background(slide, prs) -> None:
    rect = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    rect.line.fill.background()
    rect.fill.solid()
    rect.fill.fore_color.rgb = _BG
    rect.shadow.inherit = False
    # send to back
    spTree = rect._element.getparent()
    spTree.remove(rect._element)
    spTree.insert(2, rect._element)


def generate_pptx(title: str, subtitle: str | None, slides: list[dict[str, Any]]) -> str:
    """Build a .pptx and return the absolute file path.

    `slides`: list of dicts. Each dict supports keys:
        - "title": str — slide heading
        - "bullets": list[str] — bullet points (optional)
        - "body": str — single paragraph body (optional, used if no bullets)
    """
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]  # truly blank

    # --- Title slide ---
    title_slide = prs.slides.add_slide(blank)
    _paint_background(title_slide, prs)
    # accent stripe
    stripe = title_slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.6), Inches(0.18), Inches(2.3)
    )
    stripe.line.fill.background()
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = _ACCENT
    _add_textbox(title_slide, x=0.7, y=2.6, w=12, h=1.5,
                 text=title, size_pt=44, bold=True)
    if subtitle:
        _add_textbox(title_slide, x=0.7, y=4.2, w=12, h=0.8,
                     text=subtitle, size_pt=20, color=_DIM)

    # --- Content slides ---
    for s in slides:
        slide = prs.slides.add_slide(blank)
        _paint_background(slide, prs)
        slide_title = s.get("title", "")
        _add_textbox(slide, x=0.7, y=0.5, w=12, h=1.0,
                     text=slide_title, size_pt=30, bold=True)
        # accent underline
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(1.4), Inches(0.6), Inches(0.06)
        )
        underline.line.fill.background()
        underline.fill.solid()
        underline.fill.fore_color.rgb = _ACCENT

        bullets = s.get("bullets") or []
        body = s.get("body") or ""
        if bullets:
            _add_bullets(slide, x=0.9, y=1.9, w=11.5, h=5.0, items=bullets, size_pt=20)
        elif body:
            _add_textbox(slide, x=0.9, y=1.9, w=11.5, h=5.0,
                         text=body, size_pt=18, color=_INK)

    # save
    fname = f"deck-{uuid.uuid4().hex[:8]}.pptx"
    path = os.path.join(OUTPUT_DIR, fname)
    prs.save(path)
    return path
